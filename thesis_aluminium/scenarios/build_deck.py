"""Build the scenario deck on the SBTi template.

Mirrors Humphrey's steel deck from the scenarios section onward: divider, scenario
definitions, business-as-usual results, the two levers tested, comparison, questions,
then a figure annex.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

TEMPLATE = Path("/Users/parkercaswell/Downloads/Aluminium Model Deep Dive and IIASA Data Request (1).pptx")
OUTPUT = Path("/Users/parkercaswell/Desktop/Aluminium Scenario Results.pptx")
FIGURES = Path("figures")

PURPLE = RGBColor(0x4B, 0x2A, 0x6A)
DARK = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0x77, 0x77, 0x77)

DIVIDER, BLANK = 11, 19

prs = Presentation(str(TEMPLATE))
for index in range(len(prs.slides) - 1, -1, -1):
    slide_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(slide_id.rId)
    prs.slides._sldIdLst.remove(slide_id)


def add(layout_index):
    return prs.slides.add_slide(prs.slide_layouts[layout_index])


def clear_placeholders(slide):
    for shape in list(slide.placeholders):
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            shape._element.getparent().remove(shape._element)


def textbox(slide, text, left, top, width, size, bold=False, colour=DARK, italic=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.5))
    frame = box.text_frame
    frame.word_wrap = True
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    return box


def title(slide, text, size=27):
    return textbox(slide, text, 0.6, 0.35, 12.1, size, bold=True, colour=PURPLE)


def subtitle(slide, text, top=1.12, size=12):
    return textbox(slide, text, 0.6, top, 12.1, size, colour=GREY)


def note(slide, text, top=6.55, size=10):
    return textbox(slide, text, 0.6, top, 12.1, size, colour=GREY, italic=True)


def table(slide, rows, left=0.6, top=1.7, width=12.1, height=None,
          sizes=(11, 10), col_widths=None):
    n_rows, n_cols = len(rows), len(rows[0])
    height = height or min(0.33 * n_rows + 0.1, 4.8)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tbl = shape.table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(Inches(width) * w / total))
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            if c > 0:
                para.alignment = 2  # right
            for run in para.runs:
                run.font.size = Pt(sizes[0] if r == 0 else sizes[1])
                run.font.bold = r == 0
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if r == 0 else DARK
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.025)
    return tbl


def bullets(slide, items, top=1.45, left=0.7, width=11.9, size=11, gap=9):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.0))
    frame = box.text_frame
    frame.word_wrap = True
    for i, (lead, text) in enumerate(items):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        run = para.add_run(); run.text = lead
        run.font.size = Pt(size); run.font.bold = True; run.font.color.rgb = PURPLE
        run = para.add_run(); run.text = text
        run.font.size = Pt(size); run.font.color.rgb = DARK
        para.space_after = Pt(gap)
    return box


def divider(text):
    slide = add(DIVIDER)
    for shape in slide.placeholders:
        if shape.has_text_frame:
            shape.text_frame.text = text
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(36); run.font.bold = True
    return slide


def figure_slide(heading, image, caption):
    slide = add(BLANK); clear_placeholders(slide)
    title(slide, heading, size=22)
    picture = slide.shapes.add_picture(str(FIGURES / image), Inches(0.7), Inches(1.4), width=Inches(11.9))
    if picture.height > Inches(4.8):
        scale = Inches(4.8) / picture.height
        picture.height = Inches(4.8); picture.width = int(picture.width * scale)
    picture.left = int((prs.slide_width - picture.width) / 2)
    note(slide, caption, top=6.45)
    return slide


# ============================================================ 1  divider
divider("Scenarios")

# ============================================================ 2  definitions
slide = add(BLANK); clear_placeholders(slide)
title(slide, "Scenarios")
subtitle(slide, "Five runs of each plant type. Two things vary: the grid emissions assumption, and whether "
                "a plant may change its power source after its anode has been converted.")
table(slide, [
    ["Scenario", "Exogenous adjustment", "What it tests"],
    ["1. Business as usual",
     "MPP default inputs, unedited. No constraints applied.",
     "What the sector does with no intervention."],
    ["2. MPP grid, shipped switches",
     "MPP's own sector carbon budget switched on. Smelter demand corrected to the 1.5°C series so refining output balances smelting demand.",
     "Whether MPP's own 1.5°C setup reaches its own budget."],
    ["3. SBTi grid, shipped switches",
     "As above, plus grid emission factors from the SBTi power pathway V4.0.",
     "Whether a credible grid trajectory changes which technologies win."],
    ["4. MPP grid, anode unlocked",
     "As scenario 2, plus six switches added so an inert-anode plant on captive fossil power may connect to the grid, sign a power purchase agreement, or move to a reactor.",
     "Whether the missing routes in MPP's switch table are what strands captive fossil generation."],
    ["5. SBTi grid, anode unlocked",
     "Both changes together.",
     "The combined effect, and whether the two interact."],
], top=1.75, sizes=(11, 9.5), col_widths=[2.4, 5.3, 4.4], height=4.5)
note(slide, "Model is MPP commit b09472f. Code changes limited to setting the sector, a pandas 2 fix, and the demand file swap. "
            "Everything else is MPP's own data.", top=6.4)

# ============================================================ 3  BAU
slide = add(BLANK); clear_placeholders(slide)
title(slide, "Business as usual results at a glance")
table(slide, [
    ["Metric", "2025", "2050"],
    ["Smelting — fleet still on carbon anode", "100%", "68%"],
    ["Smelting — grid and PPA share", "24%", "9%"],
    ["Smelting — captive fossil share", "64%", "81%"],
    ["Smelting — emissions, Mt CO₂", "755", "1,042"],
    ["Smelting — intensity, t CO₂ / t aluminium", "11.80", "12.11"],
    ["Refining — unabated fossil boiler share", "100%", "100%"],
    ["Refining — emissions, Mt CO₂", "105", "135"],
    ["Refining — intensity, t CO₂ / t alumina", "0.85", "0.81"],
    ["Combined cumulative 2020–2050", "", "31.1 Gt CO₂"],
], top=1.45, sizes=(11, 11), col_widths=[7.0, 1.8, 1.8], width=10.6)
bullets(slide, [("Budget check   ",
    "The sector 1.5°C budget is 12.0 Gt CO₂ across both plant types. This run spends 31.1 Gt — two and a half times over. "
    "Smelting intensity rises rather than falls, because business as usual converts anodes while keeping captive coal generation.")],
    top=5.95)

# ============================================================ 4  the switch table finding
slide = add(BLANK); clear_placeholders(slide)
title(slide, "A gap in the switch table strands captive generation")
bullets(slide, [
    ("What we found   ",
     "MPP's switch table lets a plant change its power source only while it still has a carbon anode. Once the anode is "
     "converted, the power source is frozen. Inert Anode + Natural Gas and Inert Anode + Coal have exactly one destination "
     "each in the data — themselves. All three shipped variants of the file behave the same way."),
    ("Why it is a gap rather than an assumption   ",
     "Anode type and electricity supply are physically independent. A smelter that has fitted inert anodes can connect to the "
     "grid exactly as easily as one that has not. Both technologies are also classified as transition rather than end-state, "
     "so the model treats them as a stepping stone — but no step onward exists."),
    ("What it costs   ",
     "In scenario 2 this leaves 7.4 Mt of aluminium on captive gas and 1.5 Mt on captive coal in 2050, with no route out. "
     "The model does not see it coming: the ranking compares one step at a time, so it takes the anode conversion now and "
     "forecloses the power switch permanently."),
    ("What we changed   ",
     "Six switches added — Inert Anode + Natural Gas and Inert Anode + Coal, each able to reach Grid, PPA+Grid and Small "
     "Modular Reactor. Costs copied from existing rows with the same destination, verified not to vary by origin."),
], top=1.35, gap=13)

# ============================================================ 5  grid construction
slide = add(BLANK); clear_placeholders(slide)
title(slide, "Replacing the grid assumption")
subtitle(slide, "Grid emission factors from the SBTi power pathway V4.0, applied to the six grid-connected smelting technologies.")
table(slide, [
    ["Grid intensity, kg CO₂ / MWh", "2030", "2040", "2050"],
    ["MPP, China — North", "454", "227", "1"],
    ["SBTi, emerging economies", "253", "15", "1"],
    ["MPP, US", "97", "19", "0"],
    ["SBTi, advanced economies", "85", "1", "2"],
], top=1.7, sizes=(11, 11), col_widths=[5.2, 1.6, 1.6, 1.6], width=9.0)
bullets(slide, [
    ("Why it matters   ",
     "Grid-connected smelting in northern China emits 6.17 t CO₂ per tonne in 2030, against 1.51 for captive coal with capture. "
     "The model is not being perverse when it builds captive generation — it is picking the lower number."),
    ("Region mapping   ",
     "Advanced: Canada, Oceania, Rest of Europe, Scandinavia, US. Emerging: the six Chinese sub-regions, Rest of Asia, Russia, "
     "Middle East, Africa, South America. Our grouping, not an official IEA table."),
    ("Method   ",
     "Scope 2 values are rescaled by the ratio of SBTi to MPP intensity, region by region and year by year, rather than "
     "overwritten. This preserves MPP's structure, for example that a power purchase agreement is already partly renewable."),
], top=4.05)

# ============================================================ 6  which lever does what
slide = add(BLANK); clear_placeholders(slide)
title(slide, "Which change does the work")
subtitle(slide, "Smelting only, so the two levers can be read cleanly. Combined budget for smelting is 10.16 Gt.")
table(slide, [
    ["Grid assumption", "Anode switches", "2040, Mt", "2050, Mt", "Cumulative, Gt", "vs budget"],
    ["MPP", "shipped", "186", "107", "11.65", "+15%"],
    ["SBTi", "shipped", "128", "111", "10.91", "+7%"],
    ["MPP", "unlocked", "143", "70", "11.12", "+9%"],
    ["SBTi", "unlocked", "81", "24", "10.10", "−1%"],
], top=1.75, sizes=(11, 11), col_widths=[2.4, 2.2, 1.7, 1.7, 2.2, 1.7], width=10.8)
bullets(slide, [
    ("Each lever alone   ",
     "The grid assumption is worth 0.74 Gt, the added switches 0.53 Gt. Together they are worth 1.55 Gt, slightly more than "
     "the sum, and only together does smelting come inside its budget."),
    ("A caution on the grid lever   ",
     "With the SBTi grid and the switches together, 90% of world aluminium ends up grid-connected and every reactor and "
     "captive plant disappears. Nuclear goes from 13.4 Mt to zero. A result that flips that completely on one input change "
     "is close to a tie, not a robust preference. The switch fix is the more durable finding: it works on MPP's own grid "
     "and leaves a diversified fleet."),
], top=4.0)

# ============================================================ 7  comparison
slide = add(BLANK); clear_placeholders(slide)
title(slide, "Scenario comparison")
table(slide, [
    ["Metric", "Budget", "Business\nas usual", "MPP grid\nshipped", "SBTi grid\nshipped", "MPP grid\nunlocked", "SBTi grid\nunlocked"],
    ["Smelting production 2050, Mt", "—", "86.1", "67.9", "67.9", "67.7", "67.9"],
    ["Inert anode share 2050", "—", "32%", "86%", "94%", "81%", "86%"],
    ["Grid and PPA share 2050", "—", "9%", "58%", "57%", "63%", "90%"],
    ["Nuclear share 2050", "—", "0%", "18%", "18%", "20%", "0%"],
    ["Hydro share 2050", "—", "9%", "11%", "10%", "10%", "10%"],
    ["Captive fossil share 2050", "—", "81%", "13%", "15%", "6%", "0%"],
    ["Anode capture share 2050", "—", "0%", "9.7%", "5.1%", "15.5%", "2.0%"],
    ["Refining unabated boiler 2050", "—", "100%", "24%", "24%", "26%", "24%"],
    ["Smelting emissions 2050, Mt", "18", "1,042", "107", "111", "70", "24"],
    ["Refining emissions 2050, Mt", "8", "135", "26", "29", "29", "25"],
    ["Cumulative combined, Gt CO₂", "11.99", "31.12", "13.71", "12.99", "13.18", "12.14"],
    ["Against budget", "—", "+159%", "+14%", "+8%", "+10%", "+1%"],
], top=1.35, sizes=(9.5, 9.5), col_widths=[3.6, 1.3, 1.5, 1.5, 1.5, 1.5, 1.5], width=12.1, height=4.9)

# ============================================================ 8  questions
slide = add(BLANK); clear_placeholders(slide)
title(slide, "Questions for MPP")
questions = [
    "Why can a plant change its power source only while it still has a carbon anode? Inert Anode + Natural Gas and Inert Anode "
    "+ Coal have no destination other than themselves in any of the three shipped switch tables, which strands them "
    "permanently. Anode type and electricity supply are physically independent. Is this deliberate?",
    "Relatedly, a renovation is allowed once per plant for any change at all, so converting the anode spends the plant's only "
    "renovation. Should it not be once per element — one anode change and one power source change over a lifetime?",
    "The shipped switch table does not reproduce your published 1.5°C mix. Our run puts 58% of 2050 smelting on grid and power "
    "purchase agreements; your published result has 18% and puts 48% on captive fossil with capture. Which switch table and "
    "which solver version produced the published run?",
    "The refinery ranking table you ship will not load against the solver you ship — it is missing the technology lifetime and "
    "cost of capital columns. Regenerating it lets the run proceed. Was that file produced by an earlier version?",
    "Technology lifetime is zero for every grid-connected technology, and their cost of capital is 13% against 9% elsewhere. "
    "How does a zero lifetime propagate through the cost calculation?",
    "Neither of our 1.5°C runs on your shipped inputs stays inside your own carbon budget. Is the budget intended as a target "
    "the solver reaches, or only as a stopping rule?",
    "The carbon capture column in the emissions data is empty in every row, so the CO₂ storage limit has nothing to measure. "
    "Was that limit intended to be usable for aluminium?",
    "The build rate limit applies only for five years after a technology matures, and only to transition and end-state "
    "technologies. After that there is no limit at all. Intended?",
]
box = slide.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(11.9), Inches(5.6))
frame = box.text_frame; frame.word_wrap = True
for i, text in enumerate(questions, 1):
    para = frame.paragraphs[0] if i == 1 else frame.add_paragraph()
    run = para.add_run(); run.text = f"{i}.  "
    run.font.size = Pt(10); run.font.bold = True; run.font.color.rgb = PURPLE
    run = para.add_run(); run.text = text
    run.font.size = Pt(10); run.font.color.rgb = DARK
    para.space_after = Pt(7)

# ============================================================ annex
divider("Annex — figures")
figure_slide("Annual emissions against the sector carbon budget", "fig1_emissions_vs_budget.png",
             "Every 1.5°C run tracks the budget to about 2035. Only the run with both changes stays with it afterwards.")
figure_slide("Cumulative emissions, both plant types combined", "fig2_cumulative.png",
             "Cumulative emissions are what matter physically. Only the combined run comes within 1% of the budget.")
figure_slide("Smelter electricity supply", "fig3_smelter_power.png",
             "All five runs start from the same fleet. Business as usual moves further onto captive generation. "
             "Unlocking the switches removes the stranded captive fossil plants; doing that with the SBTi grid removes nuclear too.")
figure_slide("Refinery digester technology", "fig4_refinery_boiler.png",
             "Refining is almost unaffected by either change. It has minimal grid exposure and no capture route at all.")
figure_slide("Emissions intensity", "fig5_intensity.png",
             "Smelting intensity falls from 11.8 to between 0.4 and 1.6 t CO₂ per tonne depending on the scenario, and rises under business as usual.")

prs.save(str(OUTPUT))
print(f"saved {OUTPUT}")
